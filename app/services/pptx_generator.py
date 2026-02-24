"""
PPTX 생성 서비스 v2.0
- AI 생성 JSON slides[] 배열 기반으로 슬라이드 구성
- 12가지 슬라이드 타입별 전문 레이아웃
- 한국어 폰트 (맑은 고딕) 완전 지원  
- 사용자 선택 accent/bg/font 반영
- python-pptx 도형으로 비주얼 요소 구현
"""

import io
import json
import logging
from datetime import date
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

logger = logging.getLogger(__name__)


# ─── 상수 ─────────────────────────────────────────────────────────────────────
KR_FONT      = "Malgun Gothic"   # 맑은 고딕 (Windows/Office 한국어 기본)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK   = RGBColor(0x1A, 0x1A, 0x1A)
LIGHT_GRAY   = RGBColor(0xF0, 0xF2, 0xF5)
MID_GRAY     = RGBColor(0xBB, 0xBB, 0xBB)
DARK_NAVY    = RGBColor(0x0A, 0x14, 0x28)
DEFAULT_ACCENT = RGBColor(0x1E, 0x6F, 0xD9)   # 기본 포인트 컬러

# bgColor 문자열 → RGBColor 매핑
BG_COLOR_MAP = {
    "white":  RGBColor(0xFF, 0xFF, 0xFF),
    "gray":   RGBColor(0xF8, 0xF9, 0xFA),
    "cream":  RGBColor(0xFE, 0xF9, 0xEF),
    "dark":   RGBColor(0x0D, 0x11, 0x17),
    "navy":   RGBColor(0x0A, 0x14, 0x28),
}

# fontFamily 문자열 → 폰트명 매핑
FONT_MAP = {
    "gothic":      "Malgun Gothic",
    "serif":       "Batang",
    "round":       "Malgun Gothic",
    "sans-serif":  "Arial",
    "sans":        "Arial",
}

# 스타일별 기본 accent 색상 (사용자가 직접 지정한 accentColor 없을 때 폴백)
STYLE_ACCENT = {
    "mckinsey": RGBColor(0x00, 0x4F, 0x9F),
    "amazon":   RGBColor(0xFF, 0x99, 0x00),
    "ib":       RGBColor(0xC9, 0xA0, 0x3C),
    "uniflow":  RGBColor(0x7C, 0x3A, 0xED),
}


# ─── 유틸 함수 ────────────────────────────────────────────────────────────────

def _parse_hex(hex_str: str, fallback: RGBColor = DEFAULT_ACCENT) -> RGBColor:
    """#RRGGBB 헥스 문자열 → RGBColor 변환. 실패 시 fallback 반환."""
    try:
        s = str(hex_str).lstrip("#")
        if len(s) == 6:
            return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except Exception:
        pass
    return fallback


def _is_dark(color: RGBColor) -> bool:
    """배경색이 어두우면 True (밝기 기준: 0.35 이하)"""
    r, g, b = color.red / 255, color.green / 255, color.blue / 255
    return (0.299 * r + 0.587 * g + 0.114 * b) < 0.35


def _set_font_kr(run, font_name: str, size_pt: float, bold: bool,
                 color: RGBColor, italic: bool = False):
    """
    폰트 설정. 한국어를 위해 Latin + East Asian 폰트 모두 명시.
    python-pptx는 latin만 설정하면 한국어가 깨질 수 있으므로
    XML을 직접 수정해서 <a:ea> 요소도 함께 설정.
    """
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    # East Asian (한국어) 폰트 명시
    try:
        rPr = run.font._element
        ea_tag = qn("a:ea")
        ea_elem = rPr.find(ea_tag)
        if ea_elem is None:
            from lxml import etree
            ea_elem = etree.SubElement(rPr, ea_tag)
        ea_elem.set("typeface", font_name)
    except Exception:
        pass   # lxml 없거나 내부 오류 → 무시


def _add_rect(slide, left, top, width, height, fill: RGBColor,
              line: RGBColor = None, line_pt: float = 0.0):
    """색상 채운 사각형 도형 추가. 반환: shape"""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_pt) if line_pt else Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _add_oval(slide, left, top, width, height, fill: RGBColor):
    """원(타원) 도형 추가"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(9, left, top, width, height)  # 9=oval
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def _txb(slide, text: str, l, t, w, h, font, size, bold, color,
         align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """단일 텍스트 박스 추가. 반환: text_frame"""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_font_kr(run, font, size, bold, color, italic)
    return tf


def _multiline_txb(slide, lines: list, l, t, w, h, font, size, color,
                   bullet: bool = True, space_before_pt: float = 6.0):
    """여러 줄 텍스트 박스 (불릿 포함 옵션)"""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(space_before_pt)
        run = p.add_run()
        run.text = ("▸  " + line) if bullet else line
        _set_font_kr(run, font, size, False, color)
    return tf


def _add_bg(slide, prs, color: RGBColor):
    """슬라이드 전체 배경 사각형 (맨 뒤로)"""
    W, H = prs.slide_width, prs.slide_height
    bg = _add_rect(slide, 0, 0, W, H, color)
    # z-order 맨 뒤로 이동
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def _accent_bar(slide, prs, color: RGBColor, top, height=Pt(4)):
    """포인트 컬러 가로선 (제목 아래 장식)"""
    W = prs.slide_width
    bar = _add_rect(slide, Inches(0.6), top, W - Inches(1.2), height, color)
    return bar


def _deco_rect(slide, prs, color: RGBColor, alpha_color: RGBColor = None):
    """우측 하단 장식용 사각형 클러스터"""
    W, H = prs.slide_width, prs.slide_height
    _add_rect(slide, W - Inches(1.8), H - Inches(1.8), Inches(1.5), Inches(1.5),
              color if alpha_color is None else alpha_color)
    _add_rect(slide, W - Inches(1.2), H - Inches(3.0), Inches(0.5), Inches(2.0), color)


def _page_num(slide, prs, num: int, total: int, font: str, color: RGBColor):
    """우하단 페이지 번호"""
    W, H = prs.slide_width, prs.slide_height
    _txb(slide, f"{num} / {total}", W - Inches(1.5), H - Inches(0.5),
         Inches(1.3), Inches(0.4), font, 9, False, color, PP_ALIGN.RIGHT)


# ─── 슬라이드 타입별 렌더러 ─────────────────────────────────────────────────

def _render_cover(prs, slide_data: dict, palette: dict, interview_data: dict,
                  total: int):
    """
    표지: 상단 컬러 바 + 대형 제목 + 부제목 + 제안자 정보
    배경에 오른쪽 장식 사각형 추가
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height
    bg   = palette["bg"]
    acc  = palette["accent"]
    font = palette["font"]

    # 배경
    _add_bg(s, prs, bg)
    is_dark_bg = _is_dark(bg)
    title_color = WHITE if is_dark_bg else palette["title_color"]
    body_color  = RGBColor(0xCC, 0xFF, 0xFF) if is_dark_bg else NEAR_BLACK

    # 상단 포인트 컬러 바
    _add_rect(s, 0, 0, W, Inches(0.22), acc)

    # 우측 장식 사각형들
    _add_rect(s, W - Inches(3.5), 0, Inches(3.5), H, acc)
    _add_rect(s, W - Inches(3.5), 0, Inches(0.08), H,
              RGBColor(max(0, acc.red - 40), max(0, acc.green - 40), max(0, acc.blue - 40)))

    # 메인 제목 (우측 장식 영역 제외한 좌측 영역에 배치)
    title_text = (slide_data.get("title") or
                  interview_data.get("proposalTitle") or "제안서")
    _txb(s, title_text, Inches(0.6), Inches(1.5), W - Inches(4.5), Inches(2.2),
         font, 38, True, title_color, PP_ALIGN.LEFT)

    # 부제목
    subtitle = interview_data.get("proposalSubtitle") or slide_data.get("governing_message", "")
    if subtitle:
        _txb(s, subtitle, Inches(0.6), Inches(3.9), W - Inches(4.5), Inches(0.9),
             font, 17, False, MID_GRAY if not is_dark_bg else RGBColor(0xBB, 0xCC, 0xDD),
             PP_ALIGN.LEFT)

    # 수평선
    _add_rect(s, Inches(0.6), Inches(4.9), Inches(5.0), Pt(2), acc)

    # 제안자 정보
    proposer = interview_data.get("proposerInfo", "UNIFLOW")
    _txb(s, proposer, Inches(0.6), Inches(5.1), W - Inches(4.5), Inches(0.8),
         font, 13, False, body_color, PP_ALIGN.LEFT)

    # 날짜
    today_str = date.today().strftime("%Y.%m")
    _txb(s, today_str, Inches(0.6), H - Inches(0.65), Inches(3.0), Inches(0.4),
         font, 11, False, MID_GRAY, PP_ALIGN.LEFT)

    # 우측 하단(색상 바 안) 제안자 회사 약칭
    company = (interview_data.get("proposerInfo", "").split("/")[-1].strip() or "UNIFLOW")
    _txb(s, company, W - Inches(3.3), H - Inches(1.5), Inches(3.0), Inches(1.0),
         font, 18, True, WHITE, PP_ALIGN.CENTER)


def _render_executive_summary(prs, slide_data: dict, palette: dict, num: int, total: int):
    """
    핵심 요약: Governing Message + 3열 핵심 카드들
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height
    bg, acc, font = palette["bg"], palette["accent"], palette["font"]
    _add_bg(s, prs, bg)
    is_dark_bg = _is_dark(bg)
    title_c = WHITE if is_dark_bg else palette["title_color"]
    body_c  = RGBColor(0xDD, 0xEE, 0xFF) if is_dark_bg else NEAR_BLACK
    card_bg = RGBColor(0x1E, 0x2A, 0x3A) if is_dark_bg else LIGHT_GRAY

    # 슬라이드 번호 배지
    _add_rect(s, Inches(0.3), Inches(0.3), Inches(0.45), Inches(0.45), acc)
    _txb(s, str(num), Inches(0.3), Inches(0.3), Inches(0.45), Inches(0.45),
         font, 12, True, WHITE, PP_ALIGN.CENTER)

    # 제목
    _txb(s, slide_data.get("title", "핵심 요약"), Inches(0.85), Inches(0.35),
         W - Inches(1.5), Inches(0.6), font, 22, True, title_c, PP_ALIGN.LEFT)

    # Governing Message
    gm = slide_data.get("governing_message", "")
    if gm:
        _add_rect(s, Inches(0.3), Inches(0.95), W - Inches(0.6), Inches(0.65),
                  RGBColor(max(0, acc.red - 20), max(0, acc.green - 20), acc.blue), None)
        _txb(s, f"  {gm}", Inches(0.3), Inches(0.95), W - Inches(0.6), Inches(0.65),
             font, 14, True, WHITE, PP_ALIGN.LEFT, italic=True)

    # 본문을 3등분해서 카드로 배치
    talking_points = slide_data.get("talking_points", [])
    body_text = slide_data.get("body", "")
    if not talking_points and body_text:
        # 본문을 3개 포인트로 분할
        lines = [l.strip() for l in body_text.split("\n") if l.strip()]
        talking_points = lines[:3] if lines else [body_text]
    if not talking_points:
        talking_points = ["핵심 내용 1", "핵심 내용 2", "핵심 내용 3"]

    card_w   = (W - Inches(0.9)) / 3
    card_top = Inches(1.75)
    card_h   = Inches(4.0)
    nums = ["①", "②", "③"]
    for i, point in enumerate(talking_points[:3]):
        cx = Inches(0.3) + i * card_w + Pt(6) * i
        _add_rect(s, cx, card_top, card_w - Pt(6), card_h, card_bg,
                  acc, 1.5)
        # 번호 원
        _add_oval(s, cx + Inches(0.2), card_top + Inches(0.2), Inches(0.5), Inches(0.5), acc)
        _txb(s, nums[i] if i < 3 else str(i + 1),
             cx + Inches(0.2), card_top + Inches(0.2), Inches(0.5), Inches(0.5),
             font, 14, True, WHITE, PP_ALIGN.CENTER)
        _txb(s, point, cx + Inches(0.2), card_top + Inches(0.85),
             card_w - Inches(0.55), card_h - Inches(1.0),
             font, 13, False, body_c, PP_ALIGN.LEFT)

    # Visual suggestion 하단 작은 메모
    vs = slide_data.get("visual_suggestion", "")
    if vs:
        _txb(s, f"📊 {vs}", Inches(0.3), H - Inches(0.55), W - Inches(0.6), Inches(0.4),
             font, 9, False, MID_GRAY, PP_ALIGN.LEFT, italic=True)

    _page_num(s, prs, num, total, font, MID_GRAY)


def _render_content_slide(prs, slide_data: dict, palette: dict, num: int, total: int):
    """
    범용 콘텐츠 슬라이드 (problem/solution/benefit/case_study/quote/기타):
    좌측 65% 텍스트 + 우측 35% 장식 영역
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height
    bg, acc, font = palette["bg"], palette["accent"], palette["font"]
    _add_bg(s, prs, bg)
    is_dark_bg = _is_dark(bg)
    title_c = WHITE if is_dark_bg else palette["title_color"]
    body_c  = RGBColor(0xDD, 0xEE, 0xFF) if is_dark_bg else NEAR_BLACK
    accent_light = RGBColor(
        min(255, acc.red + 60 if is_dark_bg else acc.red + 180),
        min(255, acc.green + 40 if is_dark_bg else acc.green + 120),
        min(255, acc.blue + 30 if is_dark_bg else acc.blue + 80),
    )

    # 슬라이드 번호 배지
    _add_rect(s, Inches(0.3), Inches(0.3), Inches(0.45), Inches(0.45), acc)
    _txb(s, str(num), Inches(0.3), Inches(0.3), Inches(0.45), Inches(0.45),
         font, 12, True, WHITE, PP_ALIGN.CENTER)

    # 제목
    _txb(s, slide_data.get("title", ""), Inches(0.85), Inches(0.32),
         W * 0.65 - Inches(1.0), Inches(0.7), font, 22, True, title_c)

    # Governing Message 강조 바
    gm = slide_data.get("governing_message", "")
    if gm:
        gm_top = Inches(1.1)
        _add_rect(s, Inches(0.3), gm_top, W * 0.65 - Inches(0.4), Inches(0.65),
                  acc)
        _txb(s, f"  {gm}", Inches(0.3), gm_top, W * 0.65 - Inches(0.4), Inches(0.65),
             font, 14, True, WHITE, PP_ALIGN.LEFT, italic=True)
        body_top = Inches(1.85)
    else:
        _accent_bar(s, prs, acc, top=Inches(1.05))
        body_top = Inches(1.2)

    # 본문 텍스트 (좌측 65%)
    body = slide_data.get("body", "")
    body_lines = [l.strip() for l in body.split("\n") if l.strip()] if body else []
    if body_lines:
        _multiline_txb(s, body_lines, Inches(0.4), body_top,
                       W * 0.65 - Inches(0.5), H - body_top - Inches(1.4),
                       font, 14, body_c, bullet=True)

    # Talking Points (하단 태그)
    tp = slide_data.get("talking_points", [])
    if tp:
        tag_text = "  ·  ".join(tp[:4])
        _add_rect(s, Inches(0.3), H - Inches(1.2), W * 0.65 - Inches(0.4), Inches(0.9),
                  RGBColor(max(0, acc.red - 30), max(0, acc.green - 30), max(0, acc.blue - 30))
                  if is_dark_bg else accent_light)
        _txb(s, "  " + tag_text, Inches(0.3), H - Inches(1.2),
             W * 0.65 - Inches(0.4), Inches(0.9),
             font, 11, False, WHITE if is_dark_bg else acc, wrap=True)

    # 우측 35% 장식 영역
    rx = W * 0.67
    # 큰 장식 사각형
    _add_rect(s, rx, Inches(0.4), W * 0.3, H - Inches(0.8), acc)
    _add_rect(s, rx + W * 0.18, Inches(0.4), W * 0.12, H - Inches(0.8),
              RGBColor(max(0, acc.red - 50), max(0, acc.green - 50), max(0, acc.blue - 50)))
    # 장식 내 슬라이드 번호
    _txb(s, f"{num:02d}", rx + W * 0.04, H - Inches(2.0), W * 0.22, Inches(1.5),
         font, 72, True, RGBColor(0xFF, 0xFF, 0xFF) if True else acc,
         PP_ALIGN.CENTER, italic=False)
    # visual_suggestion 이탤릭
    vs = slide_data.get("visual_suggestion", "")
    if vs:
        _txb(s, vs[:40], rx + Pt(8), Inches(0.8), W * 0.29, H * 0.45,
             font, 10, False,
             RGBColor(0xFF, 0xFF, 0xFF) if True else MID_GRAY,
             PP_ALIGN.CENTER, italic=True)

    _page_num(s, prs, num, total, font, MID_GRAY if not is_dark_bg else RGBColor(0x77, 0x88, 0x99))


def _render_data_chart(prs, slide_data: dict, palette: dict, num: int, total: int):
    """
    데이터 차트 슬라이드: 제목 + Governing Message + 막대 차트 + 해석
    수치 없으면 샘플 데이터 자동 생성
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height
    bg, acc, font = palette["bg"], palette["accent"], palette["font"]
    _add_bg(s, prs, bg)
    is_dark_bg = _is_dark(bg)
    title_c = WHITE if is_dark_bg else palette["title_color"]
    body_c  = RGBColor(0xDD, 0xEE, 0xFF) if is_dark_bg else NEAR_BLACK

    # 슬라이드 번호 배지
    _add_rect(s, Inches(0.3), Inches(0.3), Inches(0.45), Inches(0.45), acc)
    _txb(s, str(num), Inches(0.3), Inches(0.3), Inches(0.45), Inches(0.45),
         font, 12, True, WHITE, PP_ALIGN.CENTER)

    # 제목
    _txb(s, slide_data.get("title", "데이터 분석"), Inches(0.85), Inches(0.32),
         W - Inches(1.5), Inches(0.65), font, 22, True, title_c)

    # Governing Message
    gm = slide_data.get("governing_message", "")
    if gm:
        _add_rect(s, Inches(0.3), Inches(1.05), W - Inches(0.6), Inches(0.55), acc)
        _txb(s, f"  {gm}", Inches(0.3), Inches(1.05), W - Inches(0.6), Inches(0.55),
             font, 13, True, WHITE, PP_ALIGN.LEFT, italic=True)

    # 차트 데이터 준비
    # body 텍스트에서 수치 추출 시도 (없으면 샘플)
    body = slide_data.get("body", "")
    chart_data_obj = ChartData()

    # 간단한 샘플 데이터 (visual_suggestion 기반)
    import re
    numbers = re.findall(r"(\d+(?:\.\d+)?)\s*%?", body)
    nums_float = [float(n) for n in numbers[:5]] if numbers else []

    if len(nums_float) >= 2:
        labels = [f"지표{i+1}" for i in range(len(nums_float))]
        chart_data_obj.categories = labels
        chart_data_obj.add_series("현황", tuple(nums_float))
    else:
        # 샘플: 도입 전/후 비교
        chart_data_obj.categories = ["도입 전", "1개월 후", "3개월 후", "6개월 후", "1년 후"]
        chart_data_obj.add_series("성과 지표", (100, 112, 128, 145, 168))

    try:
        chart_top = Inches(1.7)
        chart = s.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.5), chart_top,
            W - Inches(1.0), H - chart_top - Inches(1.3),
            chart_data_obj
        ).chart
        # 차트 색상 설정
        from pptx.util import Pt as _Pt
        plot = chart.plots[0]
        for series in plot.series:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = acc
    except Exception as e:
        logger.warning(f"[PPTX] 차트 생성 실패, 텍스트 대체: {e}")
        _multiline_txb(s, [body[:200]], Inches(0.4), Inches(1.7),
                       W - Inches(0.8), H - Inches(3.0), font, 13, body_c)

    # 차트 해석 텍스트
    tp = slide_data.get("talking_points", [])
    interp = tp[0] if tp else ""
    if interp:
        _add_rect(s, Inches(0.3), H - Inches(1.1), W - Inches(0.6), Inches(0.85), acc)
        _txb(s, f"  📌 {interp}", Inches(0.3), H - Inches(1.1),
             W - Inches(0.6), Inches(0.85), font, 12, False, WHITE, wrap=True)

    _page_num(s, prs, num, total, font, MID_GRAY)


def _render_timeline(prs, slide_data: dict, palette: dict, num: int, total: int):
    """
    타임라인 슬라이드: 가로 타임라인 (원형 마커 + 단계별 설명)
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height
    bg, acc, font = palette["bg"], palette["accent"], palette["font"]
    _add_bg(s, prs, bg)
    is_dark_bg = _is_dark(bg)
    title_c = WHITE if is_dark_bg else palette["title_color"]
    body_c  = RGBColor(0xCC, 0xDD, 0xEE) if is_dark_bg else NEAR_BLACK

    # 슬라이드 번호 배지
    _add_rect(s, Inches(0.3), Inches(0.3), Inches(0.45), Inches(0.45), acc)
    _txb(s, str(num), Inches(0.3), Inches(0.3), Inches(0.45), Inches(0.45),
         font, 12, True, WHITE, PP_ALIGN.CENTER)

    _txb(s, slide_data.get("title", "실행 계획"), Inches(0.85), Inches(0.32),
         W - Inches(1.5), Inches(0.65), font, 22, True, title_c)

    gm = slide_data.get("governing_message", "")
    if gm:
        _txb(s, gm, Inches(0.3), Inches(1.05), W - Inches(0.6), Inches(0.55),
             font, 14, True, acc, italic=True)

    # 타임라인 라인 (가로)
    line_y = Inches(3.2)
    _add_rect(s, Inches(0.5), line_y - Pt(2), W - Inches(1.0), Pt(4), acc)

    # 타임라인 포인트들
    points = slide_data.get("talking_points", [])
    body_lines = [l.strip() for l in slide_data.get("body", "").split("\n") if l.strip()]
    if not points:
        points = body_lines[:5]
    if not points:
        points = ["Phase 1", "Phase 2", "Phase 3"]

    n_pts = min(len(points), 5)
    if n_pts < 1:
        n_pts = 1
    spacing = (W - Inches(1.2)) / n_pts if n_pts > 0 else W

    for i, pt in enumerate(points[:5]):
        cx = Inches(0.5) + spacing * i + spacing * 0.5
        # 원 마커
        dot_r = Inches(0.35)
        _add_oval(s, cx - dot_r / 2, line_y - dot_r / 2, dot_r, dot_r, acc)
        _txb(s, str(i + 1), cx - dot_r / 2, line_y - dot_r / 2,
             dot_r, dot_r, font, 11, True, WHITE, PP_ALIGN.CENTER)
        # 위/아래 교대로 텍스트 배치
        if i % 2 == 0:
            # 위쪽
            _txb(s, pt[:50], cx - Inches(1.0), Inches(1.7), Inches(2.0), Inches(1.3),
                 font, 12, False, body_c, PP_ALIGN.CENTER)
        else:
            # 아래쪽
            _txb(s, pt[:50], cx - Inches(1.0), Inches(3.75), Inches(2.0), Inches(1.3),
                 font, 12, False, body_c, PP_ALIGN.CENTER)

    _page_num(s, prs, num, total, font, MID_GRAY)


def _render_comparison(prs, slide_data: dict, palette: dict, num: int, total: int):
    """
    비교 슬라이드: 2열 Before/After 또는 우리 vs 경쟁사
    우리 측 컬럼에 포인트 컬러 강조
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height
    bg, acc, font = palette["bg"], palette["accent"], palette["font"]
    _add_bg(s, prs, bg)
    is_dark_bg = _is_dark(bg)
    title_c = WHITE if is_dark_bg else palette["title_color"]
    body_c  = RGBColor(0xDD, 0xEE, 0xFF) if is_dark_bg else NEAR_BLACK
    other_bg = RGBColor(0x22, 0x2C, 0x3A) if is_dark_bg else LIGHT_GRAY

    _add_rect(s, Inches(0.3), Inches(0.3), Inches(0.45), Inches(0.45), acc)
    _txb(s, str(num), Inches(0.3), Inches(0.3), Inches(0.45), Inches(0.45),
         font, 12, True, WHITE, PP_ALIGN.CENTER)

    _txb(s, slide_data.get("title", "비교 분석"), Inches(0.85), Inches(0.32),
         W - Inches(1.5), Inches(0.65), font, 22, True, title_c)

    gm = slide_data.get("governing_message", "")
    if gm:
        _add_rect(s, Inches(0.3), Inches(1.05), W - Inches(0.6), Inches(0.55), acc)
        _txb(s, f"  {gm}", Inches(0.3), Inches(1.05), W - Inches(0.6), Inches(0.55),
             font, 13, True, WHITE, italic=True)

    col_w = (W - Inches(0.9)) / 2
    col_top = Inches(1.75)
    col_h = H - col_top - Inches(0.5)
    labels = ["기존 방식", "UNIFLOW 적용 후"]

    # 왼쪽 열 (기존)
    _add_rect(s, Inches(0.3), col_top, col_w - Pt(6), col_h, other_bg, MID_GRAY, 1.0)
    _txb(s, labels[0], Inches(0.3), col_top, col_w - Pt(6), Inches(0.6),
         font, 16, True, body_c, PP_ALIGN.CENTER)

    # 오른쪽 열 (우리 = 강조)
    _add_rect(s, Inches(0.3) + col_w + Pt(6), col_top, col_w - Pt(6), col_h, acc)
    _txb(s, labels[1], Inches(0.3) + col_w + Pt(6), col_top, col_w - Pt(6), Inches(0.6),
         font, 16, True, WHITE, PP_ALIGN.CENTER)
    # ✓ 아이콘
    _txb(s, "✓", Inches(0.3) + col_w + Pt(6) + col_w - Inches(0.6), col_top,
         Inches(0.5), Inches(0.6), font, 22, True, WHITE, PP_ALIGN.CENTER)

    body = slide_data.get("body", "")
    body_lines = [l.strip() for l in body.split("\n") if l.strip()] if body else []
    mid = len(body_lines) // 2
    left_lines  = body_lines[:mid] if body_lines else ["기존 문제점들"]
    right_lines = body_lines[mid:] if body_lines else ["개선된 결과들"]

    _multiline_txb(s, left_lines[:6], Inches(0.45), col_top + Inches(0.7),
                   col_w - Inches(0.4), col_h - Inches(0.8), font, 13, body_c)
    _multiline_txb(s, right_lines[:6], Inches(0.45) + col_w + Pt(6),
                   col_top + Inches(0.7), col_w - Inches(0.4), col_h - Inches(0.8),
                   font, 13, WHITE)

    _page_num(s, prs, num, total, font, MID_GRAY)


def _render_infographic(prs, slide_data: dict, palette: dict, num: int, total: int):
    """
    인포그래픽: 2~4개 대형 수치 가로 배열
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height
    bg, acc, font = palette["bg"], palette["accent"], palette["font"]
    _add_bg(s, prs, bg)
    is_dark_bg = _is_dark(bg)
    title_c = WHITE if is_dark_bg else palette["title_color"]
    body_c  = RGBColor(0xCC, 0xDD, 0xEE) if is_dark_bg else NEAR_BLACK

    _add_rect(s, Inches(0.3), Inches(0.3), Inches(0.45), Inches(0.45), acc)
    _txb(s, str(num), Inches(0.3), Inches(0.3), Inches(0.45), Inches(0.45),
         font, 12, True, WHITE, PP_ALIGN.CENTER)

    _txb(s, slide_data.get("title", "주요 수치"), Inches(0.85), Inches(0.32),
         W - Inches(1.5), Inches(0.65), font, 22, True, title_c)

    gm = slide_data.get("governing_message", "")
    if gm:
        _txb(s, gm, Inches(0.3), Inches(1.05), W - Inches(0.6), Inches(0.55),
             font, 14, True, acc, italic=True)

    # 수치 추출
    import re
    body = slide_data.get("body", "")
    tp   = slide_data.get("talking_points", [])

    # talking_points에서 "숫자%" 또는 "숫자배" 패턴 추출
    numbers_info = []
    sources = tp[:4] if tp else [body]
    for src in sources[:4]:
        m = re.search(r"(\d+(?:\.\d+)?)\s*(%|배|배율|점|만|억|천만|%p)?", src)
        if m:
            val = m.group(1) + (m.group(2) or "")
            label = src.replace(m.group(0), "").strip("·: ") or src
            numbers_info.append((val, label[:20]))
        else:
            numbers_info.append(("—", src[:25]))

    if not numbers_info:
        numbers_info = [("15%", "수익률 향상"), ("70%", "시간 절감"), ("95%", "고객 만족")]

    n = min(len(numbers_info), 4)
    card_w = (W - Inches(0.6)) / n
    for i, (val, label) in enumerate(numbers_info[:4]):
        cx = Inches(0.3) + i * card_w
        card_bg = acc if i == 0 else (RGBColor(0x22, 0x2C, 0x3A) if is_dark_bg else LIGHT_GRAY)
        card_tc = WHITE if i == 0 or is_dark_bg else acc
        card_lc = WHITE if i == 0 or is_dark_bg else body_c
        _add_rect(s, cx, Inches(1.7), card_w - Pt(8), H - Inches(2.2), card_bg)
        _txb(s, val, cx, Inches(2.2), card_w - Pt(8), Inches(2.0),
             font, 52, True, card_tc, PP_ALIGN.CENTER)
        _add_rect(s, cx + Inches(0.3), Inches(4.3), card_w - Inches(0.9), Pt(2),
                  card_tc)
        _txb(s, label, cx, Inches(4.5), card_w - Pt(8), Inches(1.0),
             font, 13, False, card_lc, PP_ALIGN.CENTER)

    _page_num(s, prs, num, total, font, MID_GRAY)


def _render_closing(prs, slide_data: dict, palette: dict, interview_data: dict,
                    num: int, total: int):
    """
    마무리: 커버와 유사한 디자인 + 액션 아이템 + 제안자 연락처
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height
    bg, acc, font = palette["bg"], palette["accent"], palette["font"]
    _add_bg(s, prs, bg)
    is_dark_bg = _is_dark(bg)
    title_c = WHITE if is_dark_bg else palette["title_color"]
    body_c  = RGBColor(0xCC, 0xDD, 0xEE) if is_dark_bg else NEAR_BLACK

    # 하단 포인트 컬러 바 (커버와 대칭)
    _add_rect(s, 0, H - Inches(0.22), W, Inches(0.22), acc)
    # 좌측 색깔 바
    _add_rect(s, 0, 0, Inches(3.5), H, acc)
    _add_rect(s, Inches(3.42), 0, Inches(0.08), H,
              RGBColor(max(0, acc.red - 40), max(0, acc.green - 40), max(0, acc.blue - 40)))

    # 좌측 "감사합니다" 또는 closing 제목
    closing_title = slide_data.get("title") or "감사합니다"
    _txb(s, closing_title, Inches(0.2), Inches(2.5), Inches(3.0), Inches(1.5),
         font, 28, True, WHITE, PP_ALIGN.CENTER)

    _add_rect(s, Inches(0.5), Inches(4.2), Inches(2.5), Pt(2), WHITE)

    today_str = date.today().strftime("%Y.%m")
    _txb(s, today_str, Inches(0.5), Inches(4.4), Inches(3.0), Inches(0.5),
         font, 12, False, WHITE, PP_ALIGN.CENTER)

    # 우측: governing message (다음 단계)
    gm = slide_data.get("governing_message", "")
    body = slide_data.get("body", "")
    action_text = gm or body or "다음 단계를 함께 논의해 보시겠습니까?"
    _txb(s, "다음 단계", Inches(4.0), Inches(1.2), W - Inches(4.4), Inches(0.6),
         font, 18, True, title_c)
    _add_rect(s, Inches(4.0), Inches(1.85), W - Inches(4.4), Pt(2), acc)
    _txb(s, action_text, Inches(4.0), Inches(2.0), W - Inches(4.4), Inches(2.0),
         font, 14, False, body_c, PP_ALIGN.LEFT)

    # 제안자 연락처
    proposer = interview_data.get("proposerInfo", "")
    _txb(s, "📌 연락처", Inches(4.0), Inches(4.1), W - Inches(4.4), Inches(0.5),
         font, 14, True, acc)
    _txb(s, proposer or "UNIFLOW",
         Inches(4.0), Inches(4.65), W - Inches(4.4), Inches(1.2),
         font, 14, False, body_c)

    _page_num(s, prs, num, total, font, MID_GRAY)


# ─── 타입 → 렌더러 디스패처 ──────────────────────────────────────────────────

def _dispatch_slide(prs, slide_data: dict, palette: dict, interview_data: dict,
                    num: int, total: int):
    """슬라이드 type에 따라 적합한 렌더러 호출"""
    t = str(slide_data.get("type", "")).lower()

    if t == "cover":
        _render_cover(prs, slide_data, palette, interview_data, total)
    elif t in ("executive_summary",):
        _render_executive_summary(prs, slide_data, palette, num, total)
    elif t == "data_chart":
        _render_data_chart(prs, slide_data, palette, num, total)
    elif t == "timeline":
        _render_timeline(prs, slide_data, palette, num, total)
    elif t == "comparison":
        _render_comparison(prs, slide_data, palette, num, total)
    elif t == "infographic":
        _render_infographic(prs, slide_data, palette, num, total)
    elif t == "closing":
        _render_closing(prs, slide_data, palette, interview_data, num, total)
    else:
        # problem / solution / benefit / case_study / quote / 기타 모두 범용 레이아웃
        _render_content_slide(prs, slide_data, palette, num, total)


# ─── 레거시 폴백: proposalJson 없을 때 기존 구조로 ───────────────────────────

def _legacy_generate(prs, interview_data: dict, ai_summary: Optional[str], palette: dict):
    """
    proposalJson이 없는 구형 요청 처리 (하위 호환성 유지).
    기존 고정 7장 구조를 그대로 생성.
    """
    W = prs.slide_width
    font = palette["font"]
    acc  = palette["accent"]
    bg   = palette["bg"]
    is_dark = _is_dark(bg)
    tc = WHITE if is_dark else palette["title_color"]
    bc = RGBColor(0xCC, 0xDD, 0xFF) if is_dark else NEAR_BLACK

    def legacy_slide(title: str, bullets: list):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _add_bg(s, prs, bg)
        _txb(s, title, Inches(0.6), Inches(0.35), W - Inches(1.2), Inches(0.9),
             font, 24, True, tc)
        _accent_bar(s, prs, acc, top=Inches(1.2))
        _multiline_txb(s, bullets, Inches(0.6), Inches(1.4), W - Inches(1.2), Inches(4.5),
                       font, 14, bc)

    purpose  = interview_data.get("purpose", "")
    content  = interview_data.get("coreContent", "")
    proposer = interview_data.get("proposerInfo", "UNIFLOW")

    # 표지 (레거시)
    cover_sd = {"title": interview_data.get("proposalTitle", "제안서"),
                "governing_message": "", "type": "cover"}
    _render_cover(prs, cover_sd, palette, interview_data, 7)

    # 고정 섹션
    legacy_slide("제안 배경 및 목적",
                 [purpose or "본 제안서는 비즈니스 성장 기회를 위해 작성되었습니다.",
                  "시장 변화에 선제적으로 대응하는 전략을 수립합니다."])
    legacy_slide("핵심 전략",
                 ([l.strip() for l in content.split("\n") if l.strip()][:5])
                 if content else ["핵심 가치를 중심으로 전략을 구성합니다."])
    legacy_slide("실행 계획",
                 ["Phase 1 (1~3개월): 기반 구축",
                  "Phase 2 (4~6개월): 본격 확장",
                  "Phase 3 (7~12개월): 스케일업"])
    # 기대 효과
    effects = []
    if ai_summary:
        effects = [l.strip("·-•▸ ") for l in ai_summary.split("\n") if l.strip()][:5]
    if not effects:
        effects = ["비용 절감 및 수익성 개선", "브랜드 신뢰도 강화", "지속 가능한 성장 기반 마련"]
    legacy_slide("기대 효과", effects)

    # 마무리
    closing_sd = {"title": "감사합니다", "governing_message": "함께 시작하겠습니다.",
                  "body": "", "type": "closing", "talking_points": []}
    _render_closing(prs, closing_sd, palette, interview_data, 7, 7)


# ─── 메인 함수 ────────────────────────────────────────────────────────────────

def generate_pptx(interview_data: dict, ai_summary: Optional[str] = None) -> bytes:
    """
    인터뷰 데이터 + AI 생성 proposalJson → PPTX bytes 반환.

    interview_data 키:
        proposalJson  : AI 생성 JSON 전체 (dict). 없으면 레거시 방식.
        style         : "mckinsey" | "amazon" | "ib" | "uniflow"
        layout        : "widescreen" | "a4" | "square"
        bgColor       : "white" | "dark" | "navy" | "gray" | "cream" 또는 "#RRGGBB"
        accentColor   : "#RRGGBB" 포인트 컬러
        font          : "gothic" | "serif" | "round" | "sans-serif"
        proposalTitle : 제안서 제목
        proposerInfo  : "이름 / 연락처 / 회사명"
    """
    # ── 1. 색상·폰트 팔레트 구성 ──────────────────────────────────────────
    style_key     = str(interview_data.get("style", "mckinsey")).lower()
    style_accent  = STYLE_ACCENT.get(style_key, DEFAULT_ACCENT)
    accent_raw    = str(interview_data.get("accentColor", "")).strip()
    accent        = _parse_hex(accent_raw, style_accent) if accent_raw else style_accent

    bg_raw = str(interview_data.get("bgColor", "white")).strip()
    if bg_raw.startswith("#"):
        bg = _parse_hex(bg_raw, BG_COLOR_MAP["white"])
    else:
        bg = BG_COLOR_MAP.get(bg_raw.lower(), BG_COLOR_MAP["white"])

    font_key = str(interview_data.get("font", "gothic")).lower()
    font     = FONT_MAP.get(font_key, KR_FONT)

    is_dark_bg = _is_dark(bg)
    title_color = WHITE if is_dark_bg else RGBColor(0x00, 0x20, 0x50)

    palette = {
        "bg":          bg,
        "accent":      accent,
        "font":        font,
        "title_color": title_color,
    }

    # ── 2. 프레젠테이션 크기 ──────────────────────────────────────────────
    layout_key = str(interview_data.get("layout", "widescreen")).lower()
    prs = Presentation()
    if layout_key == "a4":
        prs.slide_width  = Inches(8.27)
        prs.slide_height = Inches(11.69)
    elif layout_key == "square":
        prs.slide_width  = Inches(7.5)
        prs.slide_height = Inches(7.5)
    else:   # widescreen 16:9
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

    # ── 3. proposalJson 추출 ──────────────────────────────────────────────
    proposal = interview_data.get("proposalJson")

    # interview_data에 없으면 ai_summary에서 파싱 시도
    if proposal is None and ai_summary:
        try:
            import re as _re
            m = _re.search(r"\{[\s\S]*\}", ai_summary)
            if m:
                proposal = json.loads(m.group(0))
        except Exception:
            pass

    # ── 4. 슬라이드 생성 ─────────────────────────────────────────────────
    if proposal and isinstance(proposal.get("slides"), list) and proposal["slides"]:
        slides_list = proposal["slides"]
        total = len(slides_list)

        # 제안서 제목/부제목 interview_data에 반영 (cover 렌더러 참조용)
        if proposal.get("title") and not interview_data.get("proposalTitle"):
            interview_data["proposalTitle"] = proposal["title"]
        if proposal.get("subtitle"):
            interview_data["proposalSubtitle"] = proposal["subtitle"]

        for slide_data in slides_list:
            try:
                num = int(slide_data.get("slide_number", 0))
                _dispatch_slide(prs, slide_data, palette, interview_data, num, total)
            except Exception as e:
                logger.error(f"[PPTX] 슬라이드 {slide_data.get('slide_number')} 렌더링 오류: {e}")
                # 오류 슬라이드는 기본 텍스트 슬라이드로 대체
                safe = {
                    "title": slide_data.get("title", "슬라이드"),
                    "governing_message": slide_data.get("governing_message", ""),
                    "body": slide_data.get("body", ""),
                    "talking_points": [],
                    "visual_suggestion": "",
                    "type": "problem",
                }
                try:
                    _render_content_slide(prs, safe, palette, num, total)
                except Exception:
                    pass
    else:
        # proposalJson 없음 → 레거시 방식
        logger.warning("[PPTX] proposalJson 없음, 레거시 방식으로 생성")
        _legacy_generate(prs, interview_data, ai_summary, palette)

    # ── 5. bytes 반환 ─────────────────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
